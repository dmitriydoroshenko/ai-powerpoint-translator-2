import glob
import logging
from pptx import Presentation
from pptx.oxml import parse_xml
from wakepy import keep
from logger_config import setup_logging
from translator import translate_all
from file_utils import save_presentation
from xml_handler import XMLMetadataHandler

setup_logging()

def get_text_frames(prs):
    """Генератор, возвращающий все текстовые фреймы и их метаданные."""
    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):
            frames = []
            
            if shape.has_table:
                for r_idx, row in enumerate(shape.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        frames.append((cell.text_frame, ("table_cell", s_idx, sh_idx, r_idx, c_idx)))
            
            elif shape.has_chart:
                chart = getattr(shape, 'chart')
                if chart.has_title and chart.chart_title.has_text_frame:
                    frames.append((chart.chart_title.text_frame, ("chart_title", s_idx, sh_idx)))
                
                for axis_type in ['category_axis', 'value_axis']:
                    try:
                        axis = getattr(chart, axis_type)
                        if axis.has_title:
                            frames.append((axis.axis_title.text_frame, (f"chart_{axis_type}", s_idx, sh_idx)))
                    except (ValueError, AttributeError):
                        continue
            
            elif hasattr(shape, "text_frame") and shape.text_frame:
                frames.append((shape.text_frame, ("text_frame", s_idx, sh_idx)))

            for tf, loc in frames:
                if tf and tf.text.strip():
                    yield tf, loc

def get_chart_text_nodes(chart):
    """
    Извлекает все возможные XML-узлы с текстом из диаграммы:
    кэш данных, названия рядов и текстовые фреймы внутри серий.
    """
    nodes = []
    nodes.extend(chart._element.xpath('.//c:tx//c:v | .//c:cat//c:v'))
    rich_nodes = chart._element.xpath('.//c:tx//a:t | .//c:cat//a:t')
    nodes.extend(rich_nodes)
    
    return nodes

def process_presentation(input_file):
    """
    Выполняет полный цикл обработки PPTX: извлечение текста, 
    перевод фреймов и элементов графиков, сохранение результата.
    """
    logging.info(f"Обработка файла: {input_file}")
    try:
        prs = Presentation(input_file)
        
        items = list(get_text_frames(prs))
        text_frames, _ = zip(*items) if items else ([], [])
        
        chart_nodes = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    chart_nodes.extend(get_chart_text_nodes(getattr(shape, 'chart')))

        handlers = [XMLMetadataHandler(tf._txBody.xml) for tf in text_frames]
        chart_texts = [node.text for node in chart_nodes if node.text]
        to_translate = [h.strip() for h in handlers] + chart_texts
        
        if not to_translate:
            logging.info(f"Нет текста для перевода в {input_file}")
            return

        translated_all = translate_all(to_translate)
        offset = len(handlers)
        
        for tf, handler, t_xml in zip(text_frames, handlers, translated_all[:offset]):
            try:
                new_txBody = parse_xml(handler.restore(t_xml))
                tf._txBody.getparent().replace(tf._txBody, new_txBody)
            except Exception as e:
                logging.error(f"Ошибка восстановления фрейма: {e}")

        for node, t_text in zip(chart_nodes, translated_all[offset:]):
            if t_text:
                node.text = t_text

        save_presentation(prs, input_file)
        
    except Exception as e:
        logging.error(f"Ошибка в {input_file}: {e}")

def main():
    for input_file in glob.glob('input/*.pptx'):
        process_presentation(input_file)

if __name__ == "__main__":
    with keep.running():
        main()
